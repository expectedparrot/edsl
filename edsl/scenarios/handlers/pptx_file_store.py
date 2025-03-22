from ..file_methods import FileMethods
import os
import tempfile


class PptxMethods(FileMethods):
    suffix = "pptx"

    def extract_text(self):
        from pptx import Presentation

        self.ppt = Presentation(self.path)

        # Extract all text from slides
        full_text = []
        for slide in self.ppt.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            full_text.append("\n".join(slide_text))

        text = "\n\n".join(full_text)
        return text

    def view_system(self):
        import os
        import subprocess

        if os.path.exists(self.path):
            try:
                if (os_name := os.name) == "posix":
                    subprocess.run(["open", self.path], check=True)  # macOS
                elif os_name == "nt":
                    os.startfile(self.path)  # Windows
                else:
                    subprocess.run(["xdg-open", self.path], check=True)  # Linux
            except Exception as e:
                print(f"Error opening PPTX: {e}")
        else:
            print("PPTX file was not found.")

    def view_notebook(self):
        from pptx import Presentation
        from IPython.display import HTML, display

        prs = Presentation(self.path)

        # Create a simple HTML representation of the slides
        html_content = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(f"<p>{shape.text}</p>")

            html_content.append(
                f"""
                <div style='border: 1px solid #ccc; margin: 10px; padding: 10px;'>
                    <h3>Slide {i}</h3>
                    {''.join(slide_content)}
                </div>
            """
            )

        html = f"""
        <div style="width: 800px; height: 800px; padding: 20px; 
                   border: 1px solid #ccc; overflow-y: auto;">
            {''.join(html_content)}
        </div>
        """
        display(HTML(html))

    def example(self):
        from pptx import Presentation

        os.makedirs("test_dir", exist_ok=True)

        # Create first presentation
        ppt1 = Presentation()
        slide = ppt1.slides.add_slide(ppt1.slide_layouts[0])
        title = slide.shapes.title
        title.text = "First Presentation"
        ppt1.save("test_dir/test1.pptx")

        # Create second presentation
        ppt2 = Presentation()
        slide = ppt2.slides.add_slide(ppt2.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Second Presentation"

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            ppt2.save(tmp.name)
            tmp.close()

        return tmp.name


if __name__ == "__main__":
    import doctest
    doctest.testmod()