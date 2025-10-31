import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .warning_utils import print_info, print_success, print_error


class ReportPPTX:
    """Helper class for generating PPTX presentations from Report objects."""

    def __init__(self, report):
        """Initialize with a Report instance.

        Args:
            report: The Report instance to generate PPTX for.
        """
        self.report = report

    def generate(self, filename: str):
        """Generate a PowerPoint (.pptx) presentation.

        Each question analysis gets its own slide with visualizations and writeups.

        Args:
            filename (str): Destination pptx file path (e.g. "report.pptx").
        """
        print_info("Generating PPTX presentation…")

        prs = Presentation()

        # Add title slide
        self._add_title_slide(prs)

        # Add overview slides
        self._add_overview_slides(prs)

        # Add question summary slide
        if self.report.include_questions_table:
            self._add_question_summary_slide(prs)

        # Add individual analysis slides
        self._add_analysis_slides(prs)

        # Save the presentation
        prs.save(filename)
        print_success(f"PPTX presentation saved to {filename}")

    def _add_title_slide(self, prs):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Survey Report"
        subtitle.text = "Analysis and Insights"

        # Style the title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            37, 99, 235
        )  # Blue color

    def _add_overview_slides(self, prs):
        """Add overview slides for survey, respondents, and scenarios."""
        if self.report.include_overview:
            self._add_content_slide(prs, "Survey Overview", self.report.survey_overview)

        if self.report.include_respondents_section:
            self._add_content_slide(
                prs, "Respondent Overview", self.report.respondent_overview
            )

        if self.report.include_scenario_section:
            self._add_content_slide(
                prs, "Scenario Overview", self.report.scenario_overview
            )

    def _add_question_summary_slide(self, prs):
        """Add slide with question summary table."""
        slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout)

        # Use proper title placeholder
        slide.shapes.title.text = "Question Summary"

        # Get question data
        included_questions = set()
        for analysis in self.report.analyses:
            for question_name in analysis:
                included_questions.add(question_name)

        # Create table
        questions_data = []
        for question in self.report.results.survey.questions:
            questions_data.append(
                [
                    question.question_name,
                    question.question_type,
                    "✓" if question.question_name in included_questions else "-",
                ]
            )

        if questions_data:
            self._add_table_to_slide(
                slide,
                ["Question Name", "Question Type", "Included in Report"],
                questions_data,
                top=Inches(1.5),
            )  # Start below title

    def _add_analysis_slides(self, prs):
        """Add individual analysis slides for each question."""
        for question_names, output_dict in self.report.items():
            self._add_question_analysis_slide(prs, question_names, output_dict)

    def _add_question_analysis_slide(self, prs, question_names, output_dict):
        """Add a slide for a specific question analysis."""
        # For ThemeFinderOutput, create separate slides for each chart
        theme_finder_outputs = [
            obj
            for obj in output_dict.values()
            if obj.__class__.__name__ == "ThemeFinderOutput"
        ]

        if theme_finder_outputs:
            # Create a section divider slide for theme analysis
            self._add_theme_section_divider_slide(prs, question_names)
            # Handle theme finder with multiple slides
            for theme_output in theme_finder_outputs:
                self._add_theme_finder_slides(prs, question_names, theme_output)
        else:
            # Standard handling for other outputs - regular content slide
            slide_layout = prs.slide_layouts[5]  # Title only layout
            slide = prs.slides.add_slide(slide_layout)

            # Use proper title placeholder
            section_title = self._format_question_title(question_names)
            slide.shapes.title.text = section_title

            # Add question metadata
            current_top = self._add_question_metadata(
                slide, question_names, Inches(1.3)
            )

            # Add outputs
            for output_name, output_obj in output_dict.items():
                # Skip tables that can't be analyzed
                if (
                    hasattr(output_obj, "can_be_analyzed")
                    and not output_obj.can_be_analyzed
                ):
                    continue

                display_name = getattr(output_obj, "pretty_short_name", output_name)

                # Add output heading
                current_top = self._add_text_box(
                    slide,
                    display_name,
                    left=Inches(0.5),
                    top=current_top,
                    width=Inches(9),
                    height=Inches(0.5),
                    font_size=Pt(16),
                    bold=True,
                )

                # Add writeup if available
                current_top = self._add_writeup(
                    slide, question_names, output_name, current_top
                )

                # Add visualization
                current_top = self._add_visualization(slide, output_obj, current_top)

                # Add some spacing between outputs
                current_top += Inches(0.3)

    def _add_theme_section_divider_slide(self, prs, question_names):
        """Add a centered section divider slide for theme analysis."""
        slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout)

        # Use proper title placeholder
        section_title = self._format_question_title(question_names)
        slide.shapes.title.text = section_title

        # Add centered subtitle describing the section
        subtitle_text = "Theme Analysis"
        self._add_text_box(
            slide,
            subtitle_text,
            left=Inches(1),
            top=Inches(2.5),
            width=Inches(8),
            height=Inches(1),
            font_size=Pt(32),
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Add centered description
        description_text = "The following slides present a comprehensive analysis of themes and sentiment patterns found in the responses."
        self._add_text_box(
            slide,
            description_text,
            left=Inches(1),
            top=Inches(4),
            width=Inches(8),
            height=Inches(1.5),
            font_size=Pt(18),
            alignment=PP_ALIGN.CENTER,
        )

        # Add a decorative line or bullet points listing what's coming
        charts_preview = "• Theme Distribution\n• Sentiment Analysis by Theme\n• Individual Response Sentiment\n• Example Quotes"
        self._add_text_box(
            slide,
            charts_preview,
            left=Inches(2),
            top=Inches(5.8),
            width=Inches(6),
            height=Inches(1.5),
            font_size=Pt(16),
            alignment=PP_ALIGN.CENTER,
        )

    def _add_theme_finder_slides(self, prs, question_names, theme_output):
        """Add separate slides for each ThemeFinder chart."""
        try:
            # Get all charts that should be included
            charts = theme_output.get_all_charts_for_docx()

            if not charts:
                print_info("No charts available from ThemeFinder for PPTX")
                return

            print_info(
                f"Creating {len(charts)} separate slides for ThemeFinder charts..."
            )

            for chart_title, chart, description in charts:
                # Create a new slide for each chart
                slide_layout = prs.slide_layouts[5]  # Title only layout
                slide = prs.slides.add_slide(slide_layout)

                # Set the slide title to include both question context and chart type
                question_context = self._format_question_title(question_names)
                full_title = f"{question_context}: {chart_title}"
                slide.shapes.title.text = full_title

                # Add description
                current_top = Inches(1.3)
                if description:
                    current_top = self._add_text_box(
                        slide,
                        description,
                        left=Inches(0.5),
                        top=current_top,
                        width=Inches(9),
                        height=Inches(0.4),
                        font_size=Pt(12),
                    )
                    current_top += Inches(0.2)  # Small gap after description

                # Add the chart as a large image
                try:
                    if chart is not None:
                        if self._add_altair_as_image(
                            slide,
                            chart,
                            current_top,
                            width=Inches(9),
                            height=Inches(5.5),
                        ):
                            print_success(f"Successfully added slide: {chart_title}")
                        else:
                            print_error(f"Failed to add chart: {chart_title}")
                            self._add_text_box(
                                slide,
                                f"[Could not embed chart: {chart_title}]",
                                left=Inches(0.5),
                                top=current_top,
                                width=Inches(9),
                                height=Inches(0.5),
                                font_size=Pt(14),
                            )
                    else:
                        self._add_text_box(
                            slide,
                            f"[Chart not available: {chart_title}]",
                            left=Inches(0.5),
                            top=current_top,
                            width=Inches(9),
                            height=Inches(0.5),
                            font_size=Pt(14),
                        )
                except Exception as e:
                    print_error(f"Error adding {chart_title}: {e}")
                    self._add_text_box(
                        slide,
                        f"[Error embedding chart: {chart_title}]\nError: {str(e)}",
                        left=Inches(0.5),
                        top=current_top,
                        width=Inches(9),
                        height=Inches(1),
                        font_size=Pt(12),
                    )

        except Exception as e:
            print_error(f"Error in _add_theme_finder_slides: {e}")

    def _add_content_slide(self, prs, title: str, content: str):
        """Add a slide with title and content."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        title_shape.text = title
        content_shape.text = content

        # Style the title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)

        # Style the content
        content_shape.text_frame.paragraphs[0].font.size = Pt(18)
        content_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    def _add_slide_title(self, slide, title: str):
        """Add title to a slide using proper title placeholder."""
        try:
            # Try to use the title placeholder if it exists
            if hasattr(slide, "shapes") and slide.shapes.title:
                slide.shapes.title.text = title
                # Style the title
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
                slide.shapes.title.text_frame.paragraphs[0].font.bold = True
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                    37, 99, 235
                )
                return
        except Exception:
            pass

        # Fallback to textbox if title placeholder doesn't exist or fails
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.clear()

        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(37, 99, 235)

    def _add_text_box(
        self,
        slide,
        text: str,
        left,
        top,
        width,
        height,
        font_size=Pt(14),
        bold=False,
        alignment=PP_ALIGN.LEFT,
    ):
        """Add a text box to slide and return the bottom position."""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        p.alignment = alignment

        return top + height

    def _add_question_metadata(self, slide, question_names: list, top_position):
        """Add question metadata to slide and return new top position."""
        metadata_text = self._format_question_metadata(question_names)

        return self._add_text_box(
            slide,
            metadata_text,
            left=Inches(0.5),
            top=top_position,
            width=Inches(9),
            height=Inches(0.8),
            font_size=Pt(12),
        )

    def _add_writeup(self, slide, question_names, output_name, current_top):
        """Add writeup text if available and return new top position."""
        if (
            question_names in self.report.writeups
            and output_name in self.report.writeups[question_names]
        ):
            # Check if writeup is enabled for this analysis
            analysis_key = (
                tuple(question_names)
                if isinstance(question_names, (list, tuple))
                else (question_names,)
            )
            writeup_enabled = self.report._analysis_writeup_filters.get(
                analysis_key, True
            )
            if writeup_enabled:
                writeup_text = self.report.writeups[question_names][output_name]
                return self._add_text_box(
                    slide,
                    writeup_text,
                    left=Inches(0.5),
                    top=current_top,
                    width=Inches(9),
                    height=Inches(1),
                    font_size=Pt(12),
                )
        return current_top

    def _add_visualization(self, slide, output_obj, current_top):
        """Add visualization to slide and return new top position."""
        # Note: ThemeFinderOutput is now handled separately with individual slides
        # This method handles all other output types

        # Try to embed as image first
        if self._try_add_image(slide, output_obj, current_top):
            return current_top + Inches(3)  # Image height

        # Fallback to DataFrame as text
        if self._try_add_dataframe_text(slide, output_obj, current_top):
            return current_top + Inches(2)  # Text height

        # Final fallback
        return self._add_text_box(
            slide,
            f"[Visualization: {str(output_obj)}]",
            left=Inches(0.5),
            top=current_top,
            width=Inches(9),
            height=Inches(0.5),
            font_size=Pt(10),
        )

    def _try_add_image(self, slide, output_obj, top_position) -> bool:
        """Try to add output as image. Returns True if successful."""
        # Try SVG first (convert to PNG)
        try:
            svg_location = getattr(output_obj, "svg", None)
            if svg_location is not None:
                return self._add_svg_as_image(slide, svg_location, top_position)
        except Exception:
            pass

        # Try PNG directly
        try:
            png_location = getattr(output_obj, "png", None)
            if png_location is not None:
                return self._add_png_image(slide, png_location, top_position)
        except Exception:
            pass

        # Try Altair chart (handle all chart types)
        try:
            if hasattr(output_obj, "output"):
                chart_output = output_obj.output()
                import altair as alt

                if isinstance(
                    chart_output, (alt.Chart, alt.LayerChart, alt.FacetChart)
                ) or isinstance(chart_output, alt.TopLevelMixin):
                    return self._add_altair_as_image(slide, chart_output, top_position)
        except Exception as e:
            print_error(f"Error getting chart output for PPTX: {e}")

        return False

    def _add_svg_as_image(self, slide, svg_location, top_position) -> bool:
        """Convert SVG to PNG and add to slide."""
        try:
            svg_path = str(svg_location.path)
            import cairosvg

            # Create temporary PNG
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
                png_path = tmp_file.name

            cairosvg.svg2png(url=svg_path, write_to=png_path)

            # Add to slide
            slide.shapes.add_picture(png_path, Inches(1), top_position, width=Inches(8))

            # Clean up
            os.unlink(png_path)
            return True

        except Exception:
            return False

    def _add_png_image(self, slide, png_location, top_position) -> bool:
        """Add PNG image to slide."""
        try:
            png_path = str(png_location.path)
            slide.shapes.add_picture(png_path, Inches(1), top_position, width=Inches(8))
            return True
        except Exception:
            return False

    def _add_altair_as_image(
        self, slide, chart, top_position, width=Inches(8), height=Inches(6)
    ) -> bool:
        """Save Altair chart as PNG and add to slide."""
        try:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
                png_path = tmp_file.name

            chart.save(png_path)
            slide.shapes.add_picture(
                png_path, Inches(1), top_position, width=width, height=height
            )

            # Clean up
            os.unlink(png_path)
            return True

        except Exception:
            return False

    def _try_add_dataframe_text(self, slide, output_obj, top_position) -> bool:
        """Try to add DataFrame as formatted text."""
        try:
            if hasattr(output_obj, "output"):
                output = output_obj.output()
                import pandas as pd

                if isinstance(output, pd.DataFrame):
                    # Convert DataFrame to string representation
                    df_text = output.head(10).to_string()  # Limit to first 10 rows
                    self._add_text_box(
                        slide,
                        df_text,
                        left=Inches(0.5),
                        top=top_position,
                        width=Inches(9),
                        height=Inches(2),
                        font_size=Pt(10),
                    )
                    return True
        except Exception:
            pass
        return False

    def _add_table_to_slide(
        self, slide, headers: list, data: list, top=Inches(2), left=Inches(0.5)
    ):
        """Add a table to the slide."""
        try:
            rows = len(data) + 1  # +1 for header
            cols = len(headers)

            table_shape = slide.shapes.add_table(
                rows, cols, left, top, Inches(9), Inches(0.5 + rows * 0.3)
            )
            table = table_shape.table

            # Add headers
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)

            # Add data
            for row_idx, row_data in enumerate(data):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_data)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)

            return True
        except Exception:
            return False

    def _format_question_title(self, question_names: list) -> str:
        """Format question title for slide header."""
        if len(question_names) == 1:
            question_name = question_names[0]
            question = next(
                (
                    q
                    for q in self.report.results.survey.questions
                    if q.question_name == question_name
                ),
                None,
            )
            if question:
                question_text = question.question_text
                if len(question_text) > 60:  # Shorter for slide titles
                    question_text = question_text[:60] + "..."
                return f"{question_text} ({question_name})"
            else:
                return question_name
        else:
            return f"Analysis: {', '.join(question_names)}"

    def _format_question_metadata(self, question_names: list) -> str:
        """Format question metadata as text."""
        metadata_parts = []
        total_respondents = len(self.report.results)

        for question_name in question_names:
            question = next(
                (
                    q
                    for q in self.report.results.survey.questions
                    if q.question_name == question_name
                ),
                None,
            )
            if question:
                question_type = question.question_type
                question_options = ""
                if hasattr(question, "question_options") and question.question_options:
                    question_options = ", ".join(
                        str(opt) for opt in question.question_options[:5]
                    )  # Limit options
                    if len(question.question_options) > 5:
                        question_options += "..."
                elif hasattr(question, "option_labels") and question.option_labels:
                    question_options = ", ".join(
                        str(opt) for opt in question.option_labels[:5]
                    )
                    if len(question.option_labels) > 5:
                        question_options += "..."

                metadata = f"Question: {question_name} | Type: {question_type}"
                if question_options:
                    metadata += f" | Options: {question_options}"
                metadata += f" | Respondents: {total_respondents}"
                metadata_parts.append(metadata)

        return "\n".join(metadata_parts)
